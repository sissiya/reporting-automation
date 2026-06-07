from fastapi import FastAPI, File, UploadFile, Form, Body, BackgroundTasks, HTTPException
from fastapi.responses import HTMLResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
from typing import List, Optional
import pandas as pd
from datetime import datetime
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import re
from collections import Counter, defaultdict
import uuid
import threading

app = FastAPI(title="Extraction Application API")

BASE_DIR = os.path.dirname(__file__)
OUTPUT_DIR = os.path.join(BASE_DIR, "outputs")
FRONTEND_DIR = os.path.join(BASE_DIR, "frontend")
os.makedirs(OUTPUT_DIR, exist_ok=True)
os.makedirs(FRONTEND_DIR, exist_ok=True)

app.mount("/static", StaticFiles(directory=FRONTEND_DIR), name="static")
app.mount("/outputs", StaticFiles(directory=OUTPUT_DIR), name="outputs")

# Simple in-memory job store for async exports. For production replace with durable store.
JOBS: dict = {}

# Cached sample tables to speed up /dev_tables and summary requests
SAMPLE_TABLES = None
SAMPLE_ENTITIES = []
SAMPLE_COUNT = 0

def _load_sample_tables_once():
    global SAMPLE_TABLES, SAMPLE_ENTITIES, SAMPLE_COUNT
    sample_path = os.path.join(BASE_DIR, 'Synthèse DC - Incident.json')
    if not os.path.exists(sample_path):
        # try fallback name
        sample_path2 = os.path.join(BASE_DIR, 'sample_input.json')
        if os.path.exists(sample_path2):
            sample_path = sample_path2
        else:
            return
    try:
        import json as _json
        with open(sample_path, 'r', encoding='utf-8') as f:
            obj = _json.load(f)
        df = pd.json_normalize(obj)
        df = normalize_columns(df)
        df['Categorie'] = df.apply(categorize_row, axis=1)
        SAMPLE_TABLES = build_entity_tables(df)
        SAMPLE_ENTITIES = sorted(df['Entité'].dropna().unique().tolist()) if 'Entité' in df.columns else []
        SAMPLE_COUNT = len(df)
        print(f"[init] Loaded sample tables: {len(SAMPLE_ENTITIES)} entities, rows={SAMPLE_COUNT}")
    except Exception as e:
        print('[init] Failed loading sample tables', e)

# initialize cache at import time (fast for moderate samples)
# Do not load sample tables here; initialize later after all functions are defined.



def load_json_bytes_to_df(data_bytes: bytes) -> pd.DataFrame:
    try:
        import json
        obj = json.loads(data_bytes.decode("utf-8"))
    except Exception:
        import json
        text = data_bytes.decode("utf-8")
        if text.strip().startswith("["):
            obj = json.loads(text)
        else:
            lines = [json.loads(l) for l in text.splitlines() if l.strip()]
            obj = lines
    df = pd.json_normalize(obj)
    return df


def normalize_columns(df: pd.DataFrame) -> pd.DataFrame:
    mapping = {}
    for col in df.columns:
        normalized = re.sub(r"[^a-z0-9]", "", str(col).lower())
        if normalized in ("client", "clientname", "tenant", "societe", "company", "entite", "entité"):
            mapping[col] = "Entité"
        elif normalized in ("serveur", "server", "host", "hostname", "apparail"):
            mapping[col] = "Serveur"
        elif normalized in ("titre", "title", "subject", "objet", "description", "summary"):
            mapping[col] = "Titre"
    df = df.rename(columns=mapping)
    for target in ("Entité", "Serveur", "Titre"):
        if target not in df.columns:
            df[target] = pd.NA
    return df


def categorize_row(row) -> str:
    titre = str(row.get("Titre", "") or "")
    serveur = str(row.get("Serveur", "") or "")
    combined = f"{titre} {serveur}"
    if re.search(r"ntnx|nutanix", combined, flags=re.I):
        return "Nutanix"
    if re.search(r"cpu", titre, flags=re.I):
        return "CPU Issue"
    if re.search(r"memory|ram", titre, flags=re.I):
        return "Memory Issue"
    if re.search(r"disk|storage", titre, flags=re.I):
        return "Disk"
    if re.search(r"network|nic|link", titre, flags=re.I):
        return "Network"
    return "OS / Autres"


def build_entity_tables(df: pd.DataFrame) -> dict:
    category_order = ["CPU Issue", "Memory Issue", "Disk", "Network", "Nutanix", "OS / Autres"]
    display_map = {
        "CPU Issue": "CPU Issue",
        "Memory Issue": "Memory Issue",
        "OS / Autres": "OS service",
        "Nutanix": "Nutanix Issue",
        "Disk": "VM availability / Disk / MSSQL / autres",
        "Network": "Network",
    }
    tables = {}
    entities = sorted(df["Entité"].dropna().unique()) if "Entité" in df.columns else []
    for entity in entities:
        filtered = df[df["Entité"] == entity]
        counts = filtered["Categorie"].value_counts().reindex(category_order, fill_value=0)
        total = int(counts.sum())
        rows = []
        for key in category_order:
            display_cat = display_map.get(key, key)
            vol = int(counts.get(key, 0))
            pct = f"{int(round((vol / total) * 100, 0))}%" if total > 0 else "0%"
            rows.append({"Catégorie": display_cat, "Volume": vol, "Part du total": pct})
        tables[entity] = rows
    return tables


def build_entity_tables_from_list(obj_list: list) -> dict:
    """Build per-entity category tables from a list of raw dict objects (faster than pandas for moderate sizes)."""
    category_order = ["CPU Issue", "Memory Issue", "Disk", "Network", "Nutanix", "OS / Autres"]
    display_map = {
        "CPU Issue": "CPU Issue",
        "Memory Issue": "Memory Issue",
        "OS / Autres": "OS service",
        "Nutanix": "Nutanix Issue",
        "Disk": "VM availability / Disk / MSSQL / autres",
        "Network": "Network",
    }
    # counters per entity
    counters = defaultdict(Counter)
    entity_set = set()
    for rec in obj_list:
        # normalize keys heuristically
        ent = None
        for k in rec.keys():
            kn = re.sub(r"[^a-z0-9]", "", str(k).lower())
            if kn in ("client", "clientname", "tenant", "societe", "company", "entite", "entité"):
                ent = rec.get(k)
                break
        if not ent:
            # try common key names
            ent = rec.get('Entité') or rec.get('Entite') or rec.get('entite') or rec.get('entity')
        if not ent:
            ent = 'Unknown'
        entity_set.add(ent)
        # build tiny row dict for categorize_row
        tiny = {}
        # pick Titre-like keys
        for k in rec.keys():
            kn = re.sub(r"[^a-z0-9]", "", str(k).lower())
            if kn in ("titre", "title", "subject", "objet", "description", "summary"):
                tiny['Titre'] = rec.get(k)
                break
        # pick Serveur-like keys
        for k in rec.keys():
            kn = re.sub(r"[^a-z0-9]", "", str(k).lower())
            if kn in ("serveur", "server", "host", "hostname", "apparail"):
                tiny['Serveur'] = rec.get(k)
                break
        cat = categorize_row(tiny)
        counters[ent][cat] += 1

    tables = {}
    for ent in sorted(entity_set):
        counts = counters.get(ent, Counter())
        total = sum(counts.values())
        rows = []
        for key in category_order:
            display_cat = display_map.get(key, key)
            vol = int(counts.get(key, 0))
            pct = f"{int(round((vol / total) * 100, 0))}%" if total > 0 else "0%"
            rows.append({"Catégorie": display_cat, "Volume": vol, "Part du total": pct})
        tables[ent] = rows
    return tables



class EntitiesRequest(BaseModel):
    entities: List[str]


@app.post('/entity-summary')
async def entity_summary(payload: EntitiesRequest = Body(...)):
    """Return summary tables for the requested entities.

    Accepts JSON body {"entities": ["Entity A", "Entity B"]}.
    Optionally a file upload can be provided (multipart) - if so, use the uploaded JSON.
    """
    # Use the sample JSON on disk for server-side calculation
    sample_path = os.path.join(BASE_DIR, 'Synthèse DC - Incident.json')
    if os.path.exists(sample_path):
        import json as _json
        with open(sample_path, 'r', encoding='utf-8') as f:
            obj = _json.load(f)
        df = pd.json_normalize(obj)
    else:
        return {"error": "no data available on server"}

    df = normalize_columns(df)
    df["Categorie"] = df.apply(categorize_row, axis=1)
    tables = build_entity_tables(df)
    # filter by requested entities and validate
    result = {}
    for e in payload.entities:
        if e in tables:
            result[e] = tables[e]
        else:
            # return empty list for missing entities to keep response consistent
            result[e] = []
    return result


def create_interactive_ppt_from_df(df: pd.DataFrame, output_path: str, selected_entities: list | None = None):
    PRIMARY = RGBColor(0, 51, 102)
    ACCENT = RGBColor(0, 153, 153)
    HEADER_BG = RGBColor(23, 78, 137)
    ROW_ALT = RGBColor(240, 240, 240)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    entities = sorted(df["Entité"].dropna().unique()) if "Entité" in df.columns else []
    if selected_entities:
        # filter and preserve order of selected_entities
        entities = [e for e in selected_entities if e in entities]
    if not entities:
        raise ValueError("No entities to export")

    category_order = ["CPU Issue", "Memory Issue", "Disk", "Network", "Nutanix", "OS / Autres"]
    display_map = {
        "CPU Issue": "CPU Issue",
        "Memory Issue": "Memory Issue",
        "OS / Autres": "OS service",
        "Nutanix": "Nutanix Issue",
        "Disk": "VM availability / Disk / MSSQL / autres",
        "Network": "Network",
    }

    # Title slide
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Incident Analysis by Entity"
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = f"Total: {len(df)} incidents across {len(entities)} entities"
    except Exception:
        pass

    for entity in entities:
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
        except Exception:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
        # decorative header
        try:
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.9))
            rect.fill.solid()
            rect.fill.fore_color.rgb = ACCENT
            rect.line.fill.background()
        except Exception:
            rect = None
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(11.5), Inches(0.8))
        tf = title_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = entity
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255) if rect is not None else PRIMARY

        filtered = df[df["Entité"] == entity]
        counts = filtered["Categorie"].value_counts().reindex(category_order, fill_value=0)
        total = int(counts.sum())

        table_data = [["Catégorie", "Volume", "Part du total"]]
        for key in category_order:
            display_cat = display_map.get(key, key)
            vol = int(counts.get(key, 0))
            pct = f"{int(round((vol / total) * 100, 0))}%" if total > 0 else "0%"
            table_data.append([display_cat, str(vol), pct])

        rows = len(table_data)
        cols = len(table_data[0])
        left, top, width, height = Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.6)
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Inches(6)
        table.columns[1].width = Inches(2.2)
        table.columns[2].width = Inches(2)
        for i, row_data in enumerate(table_data):
            for j, cell_text in enumerate(row_data):
                cell = table.cell(i, j)
                cell.text = cell_text
                tfc = cell.text_frame
                tfc.word_wrap = True
                if i == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = HEADER_BG
                    para = tfc.paragraphs[0]
                    para.font.color.rgb = RGBColor(255, 255, 255)
                    para.font.bold = True
                else:
                    if i % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = ROW_ALT

        # description
        nonzero = counts[counts > 0]
        if total > 0 and not nonzero.empty:
            top_cats = nonzero.sort_values(ascending=False).head(3)
            parts = []
            for k, v in top_cats.items():
                display_cat = display_map.get(k, k)
                pct = int(round(v / total * 100, 0))
                parts.append(f"{display_cat}: {v} ({pct}% )")
            desc = f"Total incidents for {entity}: {total}. Top categories: {', '.join(parts)}."
        else:
            desc = f"No incidents recorded for {entity}."
        desc_box = slide.shapes.add_textbox(Inches(0.9), Inches(3.2), Inches(11), Inches(1.2))
        d_tf = desc_box.text_frame
        d_tf.clear()
        d_tf.text = desc
        d_tf.word_wrap = True
        d_tf.paragraphs[0].font.size = Pt(12)
        d_tf.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)

        total_box = slide.shapes.add_textbox(Inches(0.9), Inches(5.0), Inches(11), Inches(0.5))
        t_tf = total_box.text_frame
        t_tf.clear()
        t_tf.text = f"Total incidents for {entity}: {total}"
        t_tf.paragraphs[0].font.size = Pt(13)
        t_tf.paragraphs[0].font.italic = True

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    outpath = os.path.join(OUTPUT_DIR, output_path)
    prs.save(outpath)
    return outpath


def create_ppt_from_tables(tables: dict, selected_entities: list, output_path: str):
    """Create a PPT using precomputed `tables` dict where each key is an entity and value is rows of category summaries."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    PRIMARY = RGBColor(0, 51, 102)
    ACCENT = RGBColor(0, 153, 153)
    HEADER_BG = RGBColor(23, 78, 137)
    ROW_ALT = RGBColor(240, 240, 240)

    # title slide
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Incident Analysis"
    except Exception:
        pass

    for entity in selected_entities:
        rows = tables.get(entity)
        if not rows:
            continue
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
        except Exception:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
        # header
        try:
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.9))
            rect.fill.solid(); rect.fill.fore_color.rgb = ACCENT; rect.line.fill.background()
        except Exception:
            rect = None
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(11.5), Inches(0.8))
        tf = title_box.text_frame; tf.clear(); p = tf.paragraphs[0]; p.text = entity; p.font.size = Pt(22); p.font.bold=True

        # table
        rows_count = len(rows)
        cols = 3
        left, top, width, height = Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.6)
        table = slide.shapes.add_table(rows_count+1, cols, left, top, width, height).table
        table.columns[0].width = Inches(6); table.columns[1].width = Inches(2.2); table.columns[2].width = Inches(2)
        # header
        hdrs = ["Catégorie","Volume","Part du total"]
        for j,h in enumerate(hdrs):
            cell = table.cell(0,j); cell.text = h; cell.fill.solid(); cell.fill.fore_color.rgb = HEADER_BG
            para = cell.text_frame.paragraphs[0]; para.font.color.rgb = RGBColor(255,255,255); para.font.bold=True
        # rows
        for i,row in enumerate(rows, start=1):
            table.cell(i,0).text = str(row.get('Catégorie') or row.get('Categorie') or '')
            table.cell(i,1).text = str(row.get('Volume') or '')
            table.cell(i,2).text = str(row.get('Part du total') or row.get('Part_du_total') or '')

        # description
        total = sum(int(r.get('Volume') or 0) for r in rows)
        nonzero = [r for r in rows if int(r.get('Volume') or 0) > 0]
        if total>0 and nonzero:
            top_cats = sorted(nonzero, key=lambda x:int(x.get('Volume') or 0), reverse=True)[:3]
            parts = [f"{r.get('Catégorie') or r.get('Categorie')}: {r.get('Volume')}" for r in top_cats]
            desc = f"Total incidents for {entity}: {total}. Top categories: {', '.join(parts)}."
        else:
            desc = f"Total incidents for {entity}: {total}."
        desc_box = slide.shapes.add_textbox(Inches(0.9), Inches(3.2), Inches(11), Inches(1.2))
        d_tf = desc_box.text_frame; d_tf.clear(); d_tf.text = desc

    os.makedirs(OUTPUT_DIR, exist_ok=True)
    outpath = os.path.join(OUTPUT_DIR, output_path)
    prs.save(outpath)
    return outpath


def _background_create_ppt_job(file_bytes: bytes, entities_list, safe_name, jobid):
    try:
        JOBS[jobid]['status'] = 'running'
        df = load_json_bytes_to_df(file_bytes)
        df = normalize_columns(df)
        df['Categorie'] = df.apply(categorize_row, axis=1)
        tables = build_entity_tables(df)
        create_ppt_from_tables(tables, entities_list or [], safe_name)
        JOBS[jobid]['status'] = 'done'
        JOBS[jobid]['ppt_path'] = f"/outputs/{safe_name}"
    except Exception as e:
        JOBS[jobid]['status'] = 'failed'
        JOBS[jobid]['error'] = str(e)


@app.post('/generate_ppt_async')
async def generate_ppt_async(background: BackgroundTasks, file: UploadFile = File(...), entities: str = Form(None)):
    """Start background PPT generation. Returns job id which can be polled via /job-status/{jobid}."""
    data = await file.read()
    # parse entities JSON if provided
    selected = None
    if entities:
        try:
            import json as _json
            selected = _json.loads(entities) if isinstance(entities, str) else None
        except Exception:
            selected = None
    timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
    safe_name = f"report_async_{timestamp}_{uuid.uuid4().hex[:8]}.pptx"
    jobid = uuid.uuid4().hex
    JOBS[jobid] = {'status': 'pending', 'ppt_path': None, 'error': None}
    # schedule background work
    # run in a separate thread to avoid blocking event loop
    def _runner():
        _background_create_ppt_job(data, selected or [], safe_name, jobid)

    t = threading.Thread(target=_runner, daemon=True)
    t.start()
    return {'job_id': jobid, 'status': 'pending'}


@app.get('/job-status/{jobid}')
def job_status(jobid: str):
    job = JOBS.get(jobid)
    if not job:
        raise HTTPException(status_code=404, detail='job not found')
    return job



def create_simple_ppt_from_df(df: pd.DataFrame, out_path: str, title: str = None):
    prs = Presentation()
    # cover
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    if slide.shapes.title:
        slide.shapes.title.text = title or "Report"
    # add a slide per top entity if exists else one slide for whole table
    if "Entité" in df.columns:
        for ent, sub in df.groupby("Entité"):
            s = prs.slides.add_slide(prs.slide_layouts[5])
            if s.shapes.title:
                s.shapes.title.text = str(ent)
            left = Inches(0.5)
            top = Inches(1.2)
            width = Inches(9)
            height = Inches(4.5)
            txBox = s.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = f"Total incidents for {ent}: {len(sub)}"
            for i, row in sub.head(5).iterrows():
                p = tf.add_paragraph()
                p.text = " | ".join([f"{c}: {row.get(c, '')}" for c in sub.columns[:3]])
                p.level = 1
            tf.paragraphs[0].font.size = Pt(12)
    else:
        s = prs.slides.add_slide(prs.slide_layouts[5])
        if s.shapes.title:
            s.shapes.title.text = "Summary"
        left = Inches(0.5)
        top = Inches(1.2)
        width = Inches(9)
        height = Inches(4.5)
        txBox = s.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = f"Rows: {len(df)}\nColumns: {', '.join(df.columns[:10])}"

    prs.save(out_path)


@app.post("/generate_ppt_file")
async def generate_ppt_file(file: UploadFile = File(...), title: str = Form(None), entities: str = Form(None)):
    import traceback
    try:
        data = await file.read()
        df = load_json_bytes_to_df(data)
        # ensure columns are normalized and categorized before building tables
        df = normalize_columns(df)
        df["Categorie"] = df.apply(categorize_row, axis=1)
        timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
        safe_name = f"report_{timestamp}.pptx"
        out_path = os.path.join(OUTPUT_DIR, safe_name)
        # if frontend provided entities (JSON array), create PPT only for those entities
        if entities:
            try:
                import json as _json
                selected = _json.loads(entities) if isinstance(entities, str) else None
            except Exception:
                selected = None
            # build tables once and create PPT from those tables for speed
            tables = build_entity_tables(df)
            create_ppt_from_tables(tables, selected or [], safe_name)
        else:
            # default: full export using simplified PPT generator
            create_simple_ppt_from_df(df, out_path, title=title)
        return {"status": "ok", "ppt_path": f"/outputs/{safe_name}"}
    except Exception as e:
        tb = traceback.format_exc()
        print("[generate_ppt_file] error", tb)
        return {"status": "error", "error": str(e), "trace": tb}


@app.post("/upload_json")
async def upload_json(file: UploadFile = File(...)):
    data = await file.read()
    # try to parse raw JSON list and use the fast builder to compute per-entity tables
    import json as _json
    try:
        text = data.decode('utf-8')
        if text.strip().startswith('['):
            obj = _json.loads(text)
        else:
            # newline-delimited JSON
            lines = [l for l in text.splitlines() if l.strip()]
            obj = [ _json.loads(l) for l in lines ]
    except Exception:
        # fallback to pandas normalization
        df = load_json_bytes_to_df(data)
        df = normalize_columns(df)
        df["Categorie"] = df.apply(categorize_row, axis=1)
        sample = df.head(50).fillna("").to_dict(orient="records")
        entities = sorted(df["Entité"].dropna().unique().tolist()) if "Entité" in df.columns else []
        tables = build_entity_tables(df)
        # debug logging
        try:
            print(f"[upload_json] rows={len(df)}, entities={len(entities)}")
        except Exception:
            pass
        return {"rows": sample, "count": len(df), "entities": entities, "tables": tables}

    # build fast tables from obj list
    tables = build_entity_tables_from_list(obj)
    # compute simple sample/count/entities
    count = len(obj)
    # extract entities list from tables keys
    entities = sorted(list(tables.keys()))
    sample = obj[:50] if isinstance(obj, list) else []
    # debug logging for frontend issues
    try:
        print(f"[upload_json] rows={count}, entities={len(entities)}")
        print(f"[upload_json] entities sample={entities[:10]}")
        # show first entity table sizes
        if isinstance(tables, dict):
            for k in list(tables.keys())[:5]:
                print(f"[upload_json] table {k} rows={len(tables[k])}")
        else:
            print(f"[upload_json] tables type: {type(tables)}")
    except Exception:
        pass
    return {"rows": sample, "count": count, "entities": entities, "tables": tables}


@app.get('/dev_tables')
def dev_tables():
    """Return computed tables for the sample JSON file on disk for debugging."""
    # return cached sample tables when available
    if SAMPLE_TABLES is not None:
        return {"count": SAMPLE_COUNT, "entities": SAMPLE_ENTITIES, "tables": SAMPLE_TABLES}
    # fallback to on-demand build
    sample_path = os.path.join(BASE_DIR, 'Synthèse DC - Incident.json')
    if not os.path.exists(sample_path):
        return {"error": "sample file not found", "path": sample_path}
    import json
    with open(sample_path, 'r', encoding='utf-8') as f:
        obj = json.load(f)
    df = pd.json_normalize(obj)
    df = normalize_columns(df)
    df['Categorie'] = df.apply(categorize_row, axis=1)
    tables = build_entity_tables(df)
    return {"count": len(df), "entities": sorted(df['Entité'].dropna().unique().tolist()), "tables": tables}


@app.get("/", response_class=HTMLResponse)
def root():
    index_path = os.path.join(FRONTEND_DIR, "index.html")
    if os.path.exists(index_path):
        with open(index_path, "r", encoding="utf-8") as f:
            return HTMLResponse(f.read())
    return HTMLResponse("<h3>Extraction Application</h3><p>Place frontend files in /extraction_application/frontend/</p>")


if __name__ == "__main__":
    # initialize sample cache in main process to avoid NameError during import-time reloads
    try:
        _load_sample_tables_once()
    except Exception:
        pass
    import uvicorn
    uvicorn.run("server:app", host="127.0.0.1", port=8000, reload=True)
