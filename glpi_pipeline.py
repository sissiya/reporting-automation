import argparse
import os
import tempfile
from datetime import datetime
from pathlib import Path
from typing import List, Optional, Tuple

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
from matplotlib import patches
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt


def normalize_column_name(name: str) -> str:
    text = str(name).strip().lower()
    for source, target in [
        ("é", "e"),
        ("è", "e"),
        ("ê", "e"),
        ("à", "a"),
        ("â", "a"),
        ("ô", "o"),
        ("ç", "c"),
        ("ù", "u"),
        ("î", "i"),
        ("ï", "i"),
        ("'", ""),
        ('"', ""),
        (" ", "_"),
        ("-", "_"),
        ("/", "_"),
    ]:
        text = text.replace(source, target)
    while "__" in text:
        text = text.replace("__", "_")
    return text


def find_column(df: pd.DataFrame, aliases: List[str]) -> Optional[str]:
    normalized = {normalize_column_name(col): col for col in df.columns}
    for alias in aliases:
        key = normalize_column_name(alias)
        if key in normalized:
            return normalized[key]
    return None


def load_glpi_json(path: Path) -> pd.DataFrame:
    if not path.exists():
        raise FileNotFoundError(f"Fichier introuvable : {path}")

    try:
        df = pd.read_json(path, encoding="utf-8")
    except ValueError:
        df = pd.read_json(path, encoding="utf-8", lines=True)

    if df.empty:
        raise ValueError("Le fichier JSON ne contient aucune donnée valide.")
    return df


def parse_dates(df: pd.DataFrame, col_name: Optional[str]) -> pd.Series:
    if col_name is None or col_name not in df.columns:
        return pd.Series(pd.NaT, index=df.index)
    return pd.to_datetime(df[col_name], errors="coerce", dayfirst=True)


def infer_status_resolved(df: pd.DataFrame) -> pd.Series:
    status_col = find_column(df, ["statut", "etat", "status"])
    if status_col is None:
        return pd.Series(False, index=df.index)
    values = df[status_col].astype(str).fillna("")
    return values.str.contains(r"clos|resol|ferme|closed|solution", case=False, na=False)


def infer_demandeur(df: pd.DataFrame) -> pd.Series:
    demandes_col = find_column(df, ["demandeur", "requester", "utilisateur"])
    if demandes_col is None:
        return pd.Series("", index=df.index)
    return df[demandes_col].astype(str).fillna("")


def infer_priority(df: pd.DataFrame) -> pd.Series:
    priority_col = find_column(df, ["priorite", "priority", "urgence", "importance"])
    if priority_col is None:
        return pd.Series("Inconnue", index=df.index)

    raw = df[priority_col].astype(str).fillna("").str.strip().str.lower()
    mapping = {
        "critique": "P1",
        "p1": "P1",
        "1": "P1",
        "haute": "P2",
        "p2": "P2",
        "2": "P2",
        "moyenne": "P3",
        "normale": "P3",
        "basse": "P3",
        "faible": "P3",
        "p3": "P3",
        "3": "P3",
    }

    def normalize_priority(value: str) -> str:
        if not value:
            return "Inconnue"
        if value in mapping:
            return mapping[value]
        if "crit" in value:
            return "P1"
        if "haut" in value:
            return "P2"
        if "moy" in value or "norm" in value or "bas" in value or "faib" in value:
            return "P3"
        return value.upper()

    return raw.apply(normalize_priority)


def infer_assignment_group(df: pd.DataFrame) -> pd.Series:
    group_col = find_column(df, ["groupe_d_affectation", "affectation", "groupe", "assignment_group"])
    if group_col is None:
        return pd.Series("Non attribue", index=df.index)
    return df[group_col].astype(str).fillna("Non attribue")


def infer_ttr_tto(
    df: pd.DataFrame,
    created: pd.Series,
    resolved: pd.Series,
) -> Tuple[pd.Series, pd.Series]:
    ttr_col = find_column(df, ["ttr", "temps_de_resolution", "temps_resolution", "time_to_resolution"])
    tto_col = find_column(df, ["tto", "temps_de_prise_en_charge", "time_to_own", "temps_prise_en_charge"])

    def parse_duration(value: pd.Series) -> pd.Series:
        numeric = pd.to_numeric(value, errors="coerce")
        if numeric.notna().any():
            return numeric
        return pd.Series(np.nan, index=value.index)

    if ttr_col and ttr_col in df.columns:
        ttr = parse_duration(df[ttr_col])
        if ttr.isna().all():
            ttr = (resolved - created).dt.total_seconds() / 3600
    else:
        ttr = (resolved - created).dt.total_seconds() / 3600

    if tto_col and tto_col in df.columns:
        tto = parse_duration(df[tto_col])
        if tto.isna().all():
            assign_date_col = find_column(df, ["date_de_prise_en_charge", "date_prise_en_charge", "date_d_affectation", "date_affectation"])
            if assign_date_col and assign_date_col in df.columns:
                assign_date = pd.to_datetime(df[assign_date_col], errors="coerce", dayfirst=True)
                tto = (assign_date - created).dt.total_seconds() / 3600
            else:
                tto = pd.Series(np.nan, index=df.index)
    else:
        assign_date_col = find_column(df, ["date_de_prise_en_charge", "date_prise_en_charge", "date_d_affectation", "date_affectation"])
        if assign_date_col and assign_date_col in df.columns:
            assign_date = pd.to_datetime(df[assign_date_col], errors="coerce", dayfirst=True)
            tto = (assign_date - created).dt.total_seconds() / 3600
        else:
            tto = pd.Series(np.nan, index=df.index)
    return ttr, tto


def create_temp_png(fig) -> str:
    tmp_file = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
    tmp_path = tmp_file.name
    tmp_file.close()
    fig.savefig(tmp_path, dpi=150, bbox_inches="tight")
    plt.close(fig)
    return tmp_path


def create_line_chart(series_a: pd.Series, series_b: pd.Series, title: str) -> str:
    fig, ax = plt.subplots(figsize=(10, 5))
    ax.plot(series_a.index, series_a.values, marker="o", label="Tickets créés")
    ax.plot(series_b.index, series_b.values, marker="o", label="Tickets clos")
    ax.set_title(title, fontsize=14)
    ax.set_xlabel("Date")
    ax.set_ylabel("Nombre de tickets")
    ax.legend()
    ax.grid(axis="y", linestyle="--", alpha=0.4)
    fig.autofmt_xdate(rotation=45)
    return create_temp_png(fig)


def create_bar_chart(series: pd.Series, title: str, xlabel: str, ylabel: str) -> str:
    fig, ax = plt.subplots(figsize=(10, 5))
    series_sorted = series.sort_values(ascending=False)
    ax.bar(series_sorted.index.astype(str), series_sorted.values, color="#005A9C")
    ax.set_title(title, fontsize=14)
    ax.set_xlabel(xlabel)
    ax.set_ylabel(ylabel)
    ax.tick_params(axis="x", rotation=45)
    ax.grid(axis="y", linestyle="--", alpha=0.4)
    return create_temp_png(fig)


def create_gauge_chart(value: float, label: str) -> str:
    fig, ax = plt.subplots(figsize=(4, 2.5), subplot_kw={"aspect": "equal"})
    normalized = max(0.0, min(value, 100.0))
    gauge_color = "#2E7D32" if normalized >= 90 else "#C62828"
    ax.add_patch(patches.Wedge((0, 0), 1, 180, 360, facecolor="#E0E0E0", edgecolor="none"))
    ax.add_patch(patches.Wedge((0, 0), 1, 180, 180 + normalized * 1.8, facecolor=gauge_color, edgecolor="none"))
    ax.add_patch(patches.Circle((0, 0), 0.5, color="white", zorder=10))
    ax.text(0, -0.05, f"{normalized:.1f}%", fontsize=18, fontweight="bold", ha="center", va="center")
    ax.text(0, -0.4, label, fontsize=12, ha="center", va="center")
    ax.set_xlim(-1.1, 1.1)
    ax.set_ylim(-0.2, 1.1)
    ax.axis("off")
    return create_temp_png(fig)


def add_table(slide, data: pd.DataFrame, left: Cm, top: Cm, width: Cm, height: Cm) -> None:
    rows, cols = data.shape[0] + 1, data.shape[1] + 1
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    table.cell(0, 0).text = "Groupe d'affectation"
    for j, col in enumerate(data.columns, start=1):
        table.cell(0, j).text = str(col)
    for i, idx in enumerate(data.index, start=1):
        table.cell(i, 0).text = str(idx)
        for j, col in enumerate(data.columns, start=1):
            table.cell(i, j).text = str(int(data.iloc[i - 1, j - 1]))


def add_textbox(slide, left: Cm, top: Cm, width: Cm, height: Cm, text: str, size: int = 14) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(size)
    tf.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)


def build_presentation(
    df: pd.DataFrame,
    current_month: str,
    created_counts: pd.Series,
    closed_counts: pd.Series,
    group_counts: pd.Series,
    resolution_pct: dict,
    tickets_by_group_priority: pd.DataFrame,
    monthly_closed: pd.Series,
    backlog_counts: dict,
    ttr_summary: pd.DataFrame,
    sla_trend: pd.Series,
    output_path: Path,
) -> None:
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank_layout)
    title_box = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(24), Cm(1.2))
    title_tf = title_box.text_frame
    title_tf.text = f"Synthèse Mensuelle - {current_month}"
    title_tf.paragraphs[0].font.size = Pt(28)
    title_tf.paragraphs[0].font.bold = True
    title_tf.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    line_img = create_line_chart(created_counts, closed_counts, "Tickets créés vs Tickets clos")
    slide.shapes.add_picture(line_img, Cm(1), Cm(2), width=Cm(18))
    os.unlink(line_img)

    bar_img = create_bar_chart(group_counts, "Tickets par groupe d'affectation", "Groupe d'affectation", "Tickets")
    slide.shapes.add_picture(bar_img, Cm(1), Cm(12), width=Cm(18))
    os.unlink(bar_img)

    best_group = "Aucun groupe"
    if not tickets_by_group_priority.empty:
        best_group = tickets_by_group_priority.sum(axis=1).idxmax()
    add_textbox(slide, Cm(1), Cm(19.5), Cm(18), Cm(2), f"Groupe ayant résolu le plus de tickets : {best_group}.", size=16)

    # Slide 2
    slide = prs.slides.add_slide(blank_layout)
    title_box = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(24), Cm(1.2))
    title_box.text_frame.text = "Tableau de bord SLA et Priorités"
    title_box.text_frame.paragraphs[0].font.size = Pt(28)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    x_offset = Cm(1)
    for priority in ["P1", "P2", "P3"]:
        value = resolution_pct.get(priority, 0.0)
        gauge_path = create_gauge_chart(value, f"{priority} ({value:.1f}%)")
        slide.shapes.add_picture(gauge_path, x_offset, Cm(2.5), width=Cm(6))
        os.unlink(gauge_path)
        x_offset += Cm(7.5)

    add_table(slide, tickets_by_group_priority, Cm(1), Cm(12), Cm(18), Cm(6))

    # Slide 3
    slide = prs.slides.add_slide(blank_layout)
    title_box = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(24), Cm(1.2))
    title_box.text_frame.text = "Résumé historique et Backlog"
    title_box.text_frame.paragraphs[0].font.size = Pt(28)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    monthly_series = monthly_closed.reindex(pd.date_range(monthly_closed.index.min(), monthly_closed.index.max(), freq="MS"), fill_value=0) if not monthly_closed.empty else monthly_closed
    history_fig = plt.figure(figsize=(10, 5))
    ax = history_fig.add_subplot(111)
    ax.plot(monthly_series.index, monthly_series.values, marker="o", color="#005A9C")
    ax.set_title("Tendance des tickets clos", fontsize=14)
    ax.set_xlabel("Mois")
    ax.set_ylabel("Tickets clos")
    ax.grid(axis="y", linestyle="--", alpha=0.4)
    history_fig.autofmt_xdate(rotation=45)
    history_path = create_temp_png(history_fig)
    slide.shapes.add_picture(history_path, Cm(1), Cm(2), width=Cm(18))
    os.unlink(history_path)

    backlog_data = pd.DataFrame(
        [
            [">3 jours", backlog_counts.get(">3 jours", 0)],
            [">5 jours", backlog_counts.get(">5 jours", 0)],
            [">10 jours", backlog_counts.get(">10 jours", 0)],
            [">20 jours", backlog_counts.get(">20 jours", 0)],
        ],
        columns=["Tranche d'âge", "Tickets ouverts"],
    )
    backlog_table = slide.shapes.add_table(backlog_data.shape[0] + 1, backlog_data.shape[1], Cm(1), Cm(11), Cm(10), Cm(5)).table
    for col_idx, col_name in enumerate(backlog_data.columns):
        backlog_table.cell(0, col_idx).text = col_name
    for row_idx, row in backlog_data.iterrows():
        backlog_table.cell(row_idx + 1, 0).text = str(row.iloc[0])
        backlog_table.cell(row_idx + 1, 1).text = str(row.iloc[1])

    ttr_data = ttr_summary.reset_index()
    ttr_table = slide.shapes.add_table(ttr_data.shape[0] + 1, ttr_data.shape[1], Cm(12), Cm(11), Cm(10), Cm(5)).table
    for col_idx, col_name in enumerate(ttr_data.columns):
        ttr_table.cell(0, col_idx).text = str(col_name)
    for row_idx, row in ttr_data.iterrows():
        ttr_table.cell(row_idx + 1, 0).text = str(row.iloc[0])
        ttr_table.cell(row_idx + 1, 1).text = f"{row.iloc[1]:.1f}"
        ttr_table.cell(row_idx + 1, 2).text = f"{row.iloc[2]:.1f}"

    sla_img = create_bar_chart(sla_trend, "Performance SLA par mois", "Mois", "% SLA")
    slide.shapes.add_picture(sla_img, Cm(1), Cm(17), width=Cm(18))
    os.unlink(sla_img)

    prs.save(output_path)


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Transformer un export GLPI JSON en un rapport PowerPoint exécutif."
    )
    parser.add_argument("--input", default="csvjson.json", help="Chemin du fichier JSON d'entrée.")
    parser.add_argument("--output", default="SGLN_Monthly_Report.pptx", help="Chemin du fichier PowerPoint de sortie.")
    args = parser.parse_args()

    source_path = Path(args.input)
    output_path = Path(args.output)
    df = load_glpi_json(source_path)

    created_col = find_column(df, ["date_de_creation", "creation", "date_creation", "date_de_cration", "date_douverture", "date_ouverture", "date_d_ouverture", "date_d ouverture"])
    resolved_col = find_column(df, ["date_de_cloture", "date_de_resolution", "date_resolution", "date_cloture", "date_de_resolution", "date_de_resolution"])
    created = parse_dates(df, created_col)
    resolved = parse_dates(df, resolved_col)
    if created.isna().all():
        raise ValueError("Impossible de détecter la colonne de date de création dans le fichier GLPI.")

    df["Source"] = infer_demandeur(df).apply(lambda value: "Automatique" if value.strip().lower() == "monito_twin" else "Utilisateur")
    df["Priorite"] = infer_priority(df).str.upper().replace({"1": "P1", "2": "P2", "3": "P3"})
    df["Groupe d'affectation"] = infer_assignment_group(df)
    df["Date de création"] = created
    df["Date de résolution"] = resolved
    df["Résolu"] = infer_status_resolved(df)
    df["Âge"] = (pd.Timestamp.now() - df["Date de création"]).dt.days.clip(lower=0)
    df["TTR (h)"], df["TTO (h)"] = infer_ttr_tto(df, df["Date de création"], df["Date de résolution"])

    current_month = pd.Timestamp.now().strftime("%B %Y")
    current_period = pd.Timestamp.now().to_period("M")
    df_current = df[df["Date de création"].dt.to_period("M") == current_period]

    if df_current.empty:
        created_counts = pd.Series(dtype=int)
        closed_counts = pd.Series(dtype=int)
    else:
        created_counts = df_current["Date de création"].dt.date.value_counts().sort_index()
        closed_counts = df_current[df_current["Résolu"]]["Date de résolution"].dt.date.value_counts().sort_index()
        date_index = pd.date_range(df_current["Date de création"].min().date(), pd.Timestamp.now().date(), freq="D")
        created_counts = created_counts.reindex(date_index.date, fill_value=0)
        closed_counts = closed_counts.reindex(date_index.date, fill_value=0)
        created_counts.index = pd.to_datetime(created_counts.index)
        closed_counts.index = pd.to_datetime(closed_counts.index)

    group_counts = df_current["Groupe d'affectation"].value_counts()
    tickets_by_group_priority = (
        df_current[df_current["Priorite"].isin(["P1", "P2", "P3"])].groupby(["Groupe d'affectation", "Priorite"]).size().unstack(fill_value=0)
    )
    tickets_by_group_priority = tickets_by_group_priority.reindex(columns=["P1", "P2", "P3"], fill_value=0)

    resolution_pct = {}
    for priority in ["P1", "P2", "P3"]:
        priority_data = df_current[df_current["Priorite"] == priority]
        resolution_pct[priority] = (
            100.0 * priority_data[priority_data["Résolu"]].shape[0] / len(priority_data)
            if len(priority_data) > 0
            else 0.0
        )

    backlog_open = df[~df["Résolu"]].copy()
    backlog_counts = {
        ">3 jours": int((backlog_open[backlog_open["Âge"] > 3]).shape[0]),
        ">5 jours": int((backlog_open[backlog_open["Âge"] > 5]).shape[0]),
        ">10 jours": int((backlog_open[backlog_open["Âge"] > 10]).shape[0]),
        ">20 jours": int((backlog_open[backlog_open["Âge"] > 20]).shape[0]),
    }

    ttr_summary = (
        df_current[df_current["Priorite"].isin(["P1", "P2", "P3"])].groupby("Priorite")[["TTR (h)", "TTO (h)"]].mean().rename(
            columns={"TTR (h)": "Moyenne TTR (h)", "TTO (h)": "Moyenne TTO (h)"}
        )
    ).fillna(0.0)

    monthly_closed = df[df["Résolu"]].groupby(df["Date de résolution"].dt.to_period("M")).size()
    monthly_closed.index = monthly_closed.index.to_timestamp()

    sla_trend = df[df["Priorite"].isin(["P1", "P2", "P3"])].groupby(df["Date de création"].dt.to_period("M")).apply(
        lambda grp: 100.0 * grp[grp["Résolu"]].shape[0] / len(grp) if len(grp) > 0 else 0.0
    )
    sla_trend.index = sla_trend.index.to_timestamp()

    build_presentation(
        df,
        current_month,
        created_counts,
        closed_counts,
        group_counts,
        resolution_pct,
        tickets_by_group_priority,
        monthly_closed,
        backlog_counts,
        ttr_summary,
        sla_trend,
        output_path,
    )

    print(f"Rapport généré : {output_path}")


if __name__ == "__main__":
    main()
