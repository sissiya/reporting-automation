import json
import re
import os
import traceback
from typing import Any

import pandas as pd
from fastapi import FastAPI, HTTPException, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse

# PPTX
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

app = FastAPI(title="ExtractAnalyseData API")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])

# serve frontend and outputs
os.makedirs("analyse_output", exist_ok=True)
app.mount("/outputs", StaticFiles(directory="analyse_output"), name="outputs")
if os.path.isdir("frontend"):
    app.mount("/static", StaticFiles(directory="frontend"), name="static")


@app.get("/", include_in_schema=False)
async def root():
    idx = os.path.join("frontend", "index.html")
    if os.path.isfile(idx):
        return FileResponse(idx)
    return {"status": "ok", "message": "API running"}


def load_json_to_df_from_obj(obj: Any) -> pd.DataFrame:
    if isinstance(obj, dict):
        for v in obj.values():
            if isinstance(v, list) and v and isinstance(v[0], dict):
                return pd.json_normalize(v)
        return pd.json_normalize([obj])
    if isinstance(obj, list):
        return pd.json_normalize(obj)
    raise ValueError("Unsupported JSON structure")


def normalize_columns(df: pd.DataFrame) -> pd.DataFrame:
    mapping = {}
    for col in df.columns:
        normalized = re.sub(r"[^a-z0-9]", "", str(col).lower())
        if normalized in ("client", "clientname", "tenant", "societe", "company", "entite", "entité"):
            mapping[col] = "Entité"
        elif normalized in ("serveur", "server", "host", "hostname", "apparail"):
            mapping[col] = "Serveur"
        elif normalized in ("titre", "title", "subject", "objet", "description", "summary"):
            mapping[col] = "Titre"
    df = df.rename(columns=mapping)
    for target in ("Entité", "Serveur", "Titre"):
        if target not in df.columns:
            df[target] = pd.NA
    return df


def categorize_row(row) -> str:
    titre = str(row.get("Titre", "") or "")
    serveur = str(row.get("Serveur", "") or "")
    combined = f"{titre} {serveur}"
    if re.search(r"ntnx|nutanix", combined, flags=re.I):
        return "Nutanix"
    if re.search(r"cpu", titre, flags=re.I):
        return "CPU Issue"
    if re.search(r"memory|ram", titre, flags=re.I):
        return "Memory Issue"
    if re.search(r"disk|storage", titre, flags=re.I):
        return "Disk"
    if re.search(r"network|nic|link", titre, flags=re.I):
        return "Network"
    return "OS / Autres"


def analyze_json_object(obj: Any) -> dict:
    df = load_json_to_df_from_obj(obj)
    df = normalize_columns(df)
    df["Categorie"] = df.apply(categorize_row, axis=1)
    entities = sorted(df["Entité"].dropna().unique())
    category_order = ["CPU Issue", "Memory Issue", "Disk", "Network", "Nutanix", "OS / Autres"]
    results = {"total_incidents": len(df), "entities": []}
    for e in entities:
        filtered = df[df["Entité"] == e]
        counts = filtered["Categorie"].value_counts().reindex(category_order, fill_value=0).to_dict()
        results["entities"].append({"entity": e, "total": int(filtered.shape[0]), "counts": counts})
    return results


def create_interactive_ppt_from_df(df: pd.DataFrame, output_path: str = "incident_analysis_api.pptx") -> str:
    PRIMARY = RGBColor(0, 51, 102)
    ACCENT = RGBColor(0, 153, 153)
    HEADER_BG = RGBColor(23, 78, 137)
    ROW_ALT = RGBColor(240, 240, 240)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    entities = sorted(df["Entité"].dropna().unique()) if "Entité" in df.columns else []
    if not entities:
        raise ValueError("No entities to export")

    category_order = ["CPU Issue", "Memory Issue", "Disk", "Network", "Nutanix", "OS / Autres"]
    display_map = {
        "CPU Issue": "CPU Issue",
        "Memory Issue": "Memory Issue",
        "OS / Autres": "OS service",
        "Nutanix": "Nutanix Issue",
        "Disk": "VM availability / Disk / MSSQL / autres",
        "Network": "Network",
    }

    # Title slide
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Incident Analysis by Entity"
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = f"Total: {len(df)} incidents across {len(entities)} entities"
    except Exception:
        pass

    for entity in entities:
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
        except Exception:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
        # decorative header
        try:
            rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.9))
            rect.fill.solid()
            rect.fill.fore_color.rgb = ACCENT
            rect.line.fill.background()
        except Exception:
            rect = None
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(11.5), Inches(0.8))
        tf = title_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = entity
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255) if rect is not None else PRIMARY

        filtered = df[df["Entité"] == entity]
        counts = filtered["Categorie"].value_counts().reindex(category_order, fill_value=0)
        total = int(counts.sum())

        table_data = [["Catégorie", "Volume", "Part du total"]]
        for key in category_order:
            display_cat = display_map.get(key, key)
            vol = int(counts.get(key, 0))
            pct = f"{int(round((vol / total) * 100, 0))}%" if total > 0 else "0%"
            table_data.append([display_cat, str(vol), pct])

        rows = len(table_data)
        cols = len(table_data[0])
        left, top, width, height = Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.6)
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.columns[0].width = Inches(6)
        table.columns[1].width = Inches(2.2)
        table.columns[2].width = Inches(2)
        for i, row_data in enumerate(table_data):
            for j, cell_text in enumerate(row_data):
                cell = table.cell(i, j)
                cell.text = cell_text
                tfc = cell.text_frame
                tfc.word_wrap = True
                if i == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = HEADER_BG
                    para = tfc.paragraphs[0]
                    para.font.color.rgb = RGBColor(255, 255, 255)
                    para.font.bold = True
                else:
                    if i % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = ROW_ALT

        # description
        nonzero = counts[counts > 0]
        if total > 0 and not nonzero.empty:
            top_cats = nonzero.sort_values(ascending=False).head(3)
            parts = []
            for k, v in top_cats.items():
                display_cat = display_map.get(k, k)
                pct = int(round(v / total * 100, 0))
                parts.append(f"{display_cat}: {v} ({pct}% )")
            desc = f"Total incidents for {entity}: {total}. Top categories: {', '.join(parts)}."
        else:
            desc = f"No incidents recorded for {entity}."
        desc_box = slide.shapes.add_textbox(Inches(0.9), Inches(3.2), Inches(11), Inches(1.2))
        d_tf = desc_box.text_frame
        d_tf.clear()
        d_tf.text = desc
        d_tf.word_wrap = True
        d_tf.paragraphs[0].font.size = Pt(12)
        d_tf.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)

        total_box = slide.shapes.add_textbox(Inches(0.9), Inches(5.0), Inches(11), Inches(0.5))
        t_tf = total_box.text_frame
        t_tf.clear()
        t_tf.text = f"Total incidents for {entity}: {total}"
        t_tf.paragraphs[0].font.size = Pt(13)
        t_tf.paragraphs[0].font.italic = True

    # instructions slide
    try:
        instr = prs.slides.add_slide(prs.slide_layouts[1])
        instr.shapes.title.text = "How to Navigate"
        if len(instr.placeholders) > 1:
            instr.placeholders[1].text = f"Total entities: {len(entities)} | Total incidents: {len(df)}"
    except Exception:
        pass

    os.makedirs("analyse_output", exist_ok=True)
    outpath = os.path.join("analyse_output", output_path)
    prs.save(outpath)
    return outpath


# API endpoints

def load_json_file(path: str):
    with open(path, encoding="utf-8") as f:
        return json.load(f)


@app.post('/upload_json')
async def upload_json(file: UploadFile = File(...)):
    try:
        contents = await file.read()
        obj = json.loads(contents.decode('utf-8'))
        df = load_json_to_df_from_obj(obj)
        df = normalize_columns(df)
        df['Categorie'] = df.apply(categorize_row, axis=1)
        rows = df.fillna('').to_dict(orient='records')
        summary = analyze_json_object(obj)
        return {'status': 'ok', 'summary': summary, 'rows': rows}
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/analyze_file")
async def analyze_file(payload: dict):
    try:
        if "json_path" in payload:
            obj = load_json_file(payload["json_path"])
        elif "data" in payload:
            obj = payload["data"]
        else:
            raise HTTPException(status_code=400, detail="Provide 'json_path' or 'data' in the JSON body")
        summary = analyze_json_object(obj)
        return {"status": "ok", "summary": summary}
    except HTTPException:
        raise
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/generate_ppt_file")
async def generate_ppt_file(payload: dict):
    try:
        filename = payload.get("filename", "incident_analysis_api.pptx") if isinstance(payload, dict) else "incident_analysis_api.pptx"
        if "json_path" in payload:
            obj = load_json_file(payload["json_path"])
        elif "data" in payload:
            obj = payload["data"]
        else:
            raise HTTPException(status_code=400, detail="Provide 'json_path' or 'data' in the JSON body")
        df = load_json_to_df_from_obj(obj)
        df = normalize_columns(df)
        df["Categorie"] = df.apply(categorize_row, axis=1)
        out = create_interactive_ppt_from_df(df, output_path=filename)
        return {"status": "ok", "ppt_path": out}
    except HTTPException:
        raise
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


if __name__ == "__main__":
    import uvicorn

    uvicorn.run(app, host="127.0.0.1", port=8000, log_level="info")
