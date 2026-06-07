from pathlib import Path
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor

from process_and_visualize import load_and_prepare_data

DEFAULT_DATA_FILES = {
    'V1': Path('csvjson.json'),
    'V2': Path('Synthèse DC - Incident.json'),
}


def get_summary_data(suffix, data_files):
    data_file = data_files.get(suffix, DEFAULT_DATA_FILES.get(suffix))
    if not data_file or not data_file.exists():
        return None

    df = load_and_prepare_data(data_file, label=f"caption-{suffix}")
    total = len(df)
    auto_count = int((df['Source'] == 'Automatique').sum()) if 'Source' in df.columns else 0
    resolved_count = int(df['Date_Resolution'].notna().sum()) if 'Date_Resolution' in df.columns else int(df['Est_Clos'].sum() if 'Est_Clos' in df.columns else 0)
    closed_pct = resolved_count / total * 100 if total else 0
    top_group = None
    top_group_count = 0
    if 'Groupe' in df.columns:
        group_counts = df['Groupe'].value_counts()
        if len(group_counts) > 0:
            top_group = group_counts.index[0]
            top_group_count = int(group_counts.iloc[0])

    sla_scores = {}
    for code in ['P1', 'P2', 'P3']:
        sub = df[df['Priorite_Code'] == code] if 'Priorite_Code' in df.columns else df.iloc[0:0]
        total_code = len(sub)
        resolved_on_time = int((sub['SLA_Depasse'].astype(str) == 'Non').sum()) if total_code else 0
        sla_scores[code] = {
            'count': total_code,
            'score': resolved_on_time / total_code * 100 if total_code else None,
            'resolved_on_time': resolved_on_time,
        }

    backlog_df = df[~df['Est_Clos']] if 'Est_Clos' in df.columns else df.iloc[0:0]
    backlog_total = len(backlog_df)
    backlog_over_10 = int((backlog_df['Age_Jours'] > 10).sum()) if 'Age_Jours' in backlog_df.columns else 0

    return {
        'total': total,
        'auto_count': auto_count,
        'auto_pct': auto_count / total * 100 if total else 0,
        'resolved_count': resolved_count,
        'closed_pct': closed_pct,
        'top_group': top_group,
        'top_group_count': top_group_count,
        'sla_scores': sla_scores,
        'backlog_total': backlog_total,
        'backlog_over_10': backlog_over_10,
    }


def get_slide_caption_lines(image_name, summary):
    if summary is None:
        return [
            "Analyse disponible si les fichiers JSON de données sont présents.",
        ]

    if image_name == 'slide1_courbe_groupes':
        return [
            f"Ce graphique compare {summary['total']} incidents créés par rapport aux {summary['resolved_count']} incidents résolus ce mois-ci.",
            "L'alignement des courbes montre que l'équipe absorbe la charge en temps réel.",
        ]
    if image_name == 'slide1_source':
        return [
            f"Sur {summary['total']} tickets, {summary['auto_pct']:.0f}% proviennent de sources automatiques (monito).",
            "La dominance des sources automatiques confirme que la supervision détecte les incidents avant les signalements utilisateurs.",
        ]
    if image_name == 'slide2_jauges':
        lines = []
        for code in ['P1', 'P2', 'P3']:
            entry = summary['sla_scores'].get(code, {})
            if entry.get('score') is None:
                lines.append(f"SLA {code} : aucun ticket {code} ce mois-ci.")
            else:
                lines.append(f"SLA {code} à {entry['score']:.1f}% : Performance conforme aux objectifs.")
        lines.append("Ce score mesure le respect de nos engagements de délais (SLA) et reflète l'efficacité de l'équipe.")
        return lines
    if image_name == 'slide2_tableau':
        if summary['top_group']:
            return [
                f"La table couvre {summary['total']} tickets et montre que {summary['top_group']} est le groupe le plus sollicité ({summary['top_group_count']} tickets).",
                "Ce diagnostic aide à prioriser les actions sur les équipes les plus impactées.",
            ]
        return [
            f"La table couvre {summary['total']} tickets répartis par priorité et groupe.",
            "Elle aide à détecter les secteurs qui nécessitent un rééquilibrage de charge.",
        ]
    if image_name == 'slide3_historique':
        return [
            f"L'historique montre {summary['total']} tickets créés et {summary['resolved_count']} résolus ce mois-ci.",
            "Cette tendance permet d'anticiper les pics et d'ajuster les priorités opérationnelles.",
        ]
    if image_name == 'slide3_backlog_ttr':
        return [
            f"Le backlog contient {summary['backlog_total']} tickets ouverts, dont {summary['backlog_over_10']} de plus de 10 jours.",
            "Un backlog sans tickets de plus de 10 jours garantit la stabilité du service.",
        ]

    return [
        "Analyse opérationnelle disponible.",
    ]


def add_title(slide, text: str) -> None:
    title_box = slide.shapes.add_textbox(Cm(1), Cm(0.5), Cm(24), Cm(2))
    tf = title_box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = 1  # center


def add_caption(slide, left_cm, top_cm, width_cm, lines):
    caption_box = slide.shapes.add_textbox(left_cm, top_cm, width_cm, Cm(3))
    tf = caption_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Cm(0.1)
    tf.margin_right = Cm(0.1)
    tf.margin_top = Cm(0.05)
    tf.margin_bottom = Cm(0.05)
    tf.clear()
    for idx, line in enumerate(lines):
        if idx == 0:
            p = tf.paragraphs[0]
            p.text = line
        else:
            p = tf.add_paragraph()
            p.text = line
        p.font.size = Pt(10)
        p.font.name = 'Calibri'
        p.font.color.rgb = RGBColor(30, 30, 30)
        p.space_after = Pt(2)


def add_images_with_captions(slide, image_paths, summary, top_cm=2.8, spacing_cm=1.0) -> None:
    slide_width = Cm(25.4)
    image_width = (slide_width - Cm(2.4) - Cm(spacing_cm)) / 2
    image_height = Cm(7.5)
    left_positions = [Cm(1.2), Cm(1.2) + image_width + Cm(spacing_cm)]

    for idx, image_path in enumerate(image_paths[:2]):
        left = left_positions[idx]
        slide.shapes.add_picture(str(image_path), left, Cm(top_cm), width=image_width, height=image_height)
        image_name = image_path.stem.rsplit('_', 1)[0]
        caption_lines = get_slide_caption_lines(image_name, summary)
        caption_top = Cm(top_cm) + image_height + Cm(0.35)
        add_caption(slide, left, caption_top, image_width, caption_lines)


def build_report(output_path: Path, images_dir: Path, suffixes=('V1', 'V2'), data_files=None) -> None:
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    data_files = data_files or DEFAULT_DATA_FILES

    slide_pairs = [
        ('slide1_courbe_groupes', 'slide1_source'),
        ('slide2_jauges', 'slide2_tableau'),
        ('slide3_historique', 'slide3_backlog_ttr'),
    ]

    first_slide = True
    for suffix in suffixes:
        summary = get_summary_data(suffix, data_files)
        for image_names in slide_pairs:
            slide = prs.slides.add_slide(blank_layout)
            if first_slide:
                add_title(slide, 'Rapport Mensuel GLPI — V1 & V2')
                first_slide = False
            add_images_with_captions(slide, [images_dir / f"{name}_{suffix}.png" for name in image_names], summary)

    prs.save(output_path)


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Créer un PowerPoint à partir de fichiers PNG de visualisation GLPI.")
    parser.add_argument("--images-dir", default="slides", help="Dossier contenant les fichiers PNG.")
    parser.add_argument("--output", default="SGLN_Monthly_Report_V1_V2.pptx", help="Chemin du fichier PowerPoint de sortie.")
    parser.add_argument("--suffixes", nargs='+', default=['V1', 'V2'], help="Suffixes de source à inclure dans le rapport.")
    parser.add_argument("--data-file-v1", default="csvjson.json", help="Fichier JSON source pour V1.")
    parser.add_argument("--data-file-v2", default="Synthèse DC - Incident.json", help="Fichier JSON source pour V2.")
    args = parser.parse_args()

    images_dir = Path(args.images_dir)
    output_file = Path(args.output)
    data_files = {
        'V1': Path(args.data_file_v1),
        'V2': Path(args.data_file_v2),
    }
    build_report(output_file, images_dir, suffixes=args.suffixes, data_files=data_files)
    print(f"PowerPoint créé : {output_file}")